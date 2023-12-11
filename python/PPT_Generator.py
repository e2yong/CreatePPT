from Crawler_for_ppt import *
from Text_for_ppt import *

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Pt, Inches
import setting

RESULT_DIR = setting.RESULT_DIR

### 사용자 지정 옵션
# user_data: PPT로 만들 텍스트 파일명(전처리 모듈의 결과물)
# last_text: 마지막 마무리 페이지에 넣을 말
# font_name: 적용할 폰트(해당 폰트가 시스템에 저장되어 있어야 함)
# template_dir: 템플릿 저장 경로(템플릿 사용)
# template_chk: 템플릿 사용 여부(0 or 1, 0이면 기본 템플릿)
last_text = "Q&A"
font_name = "NanumGothic"
template_dir = os.getcwd() + '/' + 'template_science'
template_chk = 0

# 레이아웃 타입
l_type = {"title_slide": 0, 
        "title_and_content": 1, 
        "section_header": 2, 
        "two_content": 3, 
        "comparison": 4, 
        "title_only": 5,
        "blank": 6, 
        "content_with_caption": 7, 
        "picture_with_caption": 8
        }

# 슬라이드 타입 판별
def search_for_slide_type(text):
    tags = ["[L_TS]", "[L_CS]", "[L_IS]", "[L_THS]"]
    found_text = next((s for s in tags if s in text), None)
    return found_text

# 폰트 조정(제목)
def title_font(slide, font_name, font_size, is_bold=None):
    for paragraph in slide.shapes.title.text_frame.paragraphs:
        paragraph.font.name = font_name
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = is_bold

# 폰트 조정(본문)
def text_font(slide, idx, font_name, font_size, is_bold=None):
    for paragraph in slide.placeholders[idx].text_frame.paragraphs:
        paragraph.font.name = font_name
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = is_bold

# 폰트 크기가 가장 큰 텍스트 변경
def change_biggest_font(slide, title):
    largest_font_size = 0
    title_chk = 0

    def find_largest_font_size(slide):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size is not None:
                            font_size = run.font.size.pt
                            if font_size > largest_font_size:
                                largest_font_size = font_size
                find_largest_font_text(slide, largest_font_size)
            else:
                return title_chk
        
    def find_largest_font_text(slide, largest_font_size):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size.pt == largest_font_size:
                            old_length = len(run.text)
                            run.text = title
                            target_text = run.text
                            target_shape = shape
        
        if len(target_text) > old_length:
            diff = int(target_shape.width.pt) / len(target_text)
            target_shape.width += Pt(diff)

    find_largest_font_size(slide)

    return title_chk

# 위아래 선 긋기
def insert_two_lines(slide, shape_width, shape_height):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, slide.shapes[0].left, slide.shapes[0].top, shape_width, 4)
    line.fill.background()
    line.line.color.rgb = RGBColor(0, 0, 0)
    line.line.width = Pt(2)

    line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, slide.shapes[0].left, slide.shapes[0].top + shape_height, shape_width, 4)
    line2.fill.background()
    line2.line.color.rgb = RGBColor(0, 0, 0)
    line2.line.width = Pt(2)

### PPT 생성 모듈 ###
def generate_ppt(extract_txt):
    if template_chk:
        dir_path = template_dir
        files_in_dir = os.listdir(dir_path)
        files = [file for file in files_in_dir if os.path.isfile(os.path.join(dir_path, file))]
        template = files[random.randint(0, len(files) - 1)]
        prs = Presentation(template)
    else:
        prs = Presentation("theme0.pptx")

    # PPT를 저장하는 함수
    def save_ppt():
        ppt_name = re.sub(r"[^\uAC00-\uD7A30-9a-zA-Z\s]", "", find_title())
        prs.save(f"{RESULT_DIR}/{ppt_name}.pptx")

    ### 슬라이드 제작 함수 ###
    # 모든 슬라이드 삭제
    def delete_all_slides():
        for i in range(len(prs.slides) - 1, -1, -1):
            r_id = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(r_id)
            del prs.slides._sldIdLst[i]

    # 마지막 슬라이드
    def create_last_slide(last_text):
        layout = prs.slide_layouts[l_type["title_only"]]
        last_slide = prs.slides.add_slide(layout)

        width = Inches(6)
        height = Inches(2)

        last_slide.shapes.title.text = last_text
        last_slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        last_slide.shapes[0].width = width
        last_slide.shapes[0].height = height

        last_slide.shapes[0].top = Pt((prs.slide_height.pt - last_slide.shapes[0].height.pt) / 2)
        last_slide.shapes[0].left = Pt((prs.slide_width.pt - last_slide.shapes[0].width.pt) / 2)

        title_font(last_slide, font_name, 44, 1)

        insert_two_lines(last_slide, last_slide.shapes[0].width, last_slide.shapes[0].height)

    # 타이틀 슬라이드
    def create_title_slide(title, subtitle):
        title_slide = prs.slides[0]

        if not change_biggest_font(title_slide, title):
            title_slide.shapes.title.text = title
            title_slide.placeholders[1].text = subtitle

        title_font(title_slide, font_name, 56, 1)
        text_font(title_slide, 1, font_name, 30)
        insert_two_lines(title_slide, title_slide.shapes[0].width, title_slide.shapes[0].top + title_slide.shapes[0].height)

    # 목차
    def create_index(lines):
        pattern = re.compile(r'\[TITLE\](.*?)\[/TITLE\]')
        matches = pattern.findall(''.join(lines))
        contents = matches[1:-1]
        for i in range(0,len(contents)):
            contents[i] + '\n'

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.text = "Contents" + '\n'
        p.font.size = Pt(36)

        for content in contents:
            p = text_frame.add_paragraph()
            p.text = content
            p.space_after = Inches(0.2)

    # 섹션 헤더 슬라이드
    def create_section_header_slide(title):
        layout = prs.slide_layouts[l_type["section_header"]]

        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        slide.shapes.title.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        width = Inches(6)
        height = Inches(2)

        slide.shapes[0].width = width
        slide.shapes[0].height = height

        slide.shapes[0].top = Pt((prs.slide_height.pt - slide.shapes[0].height.pt) / 2)
        slide.shapes[0].left = Pt((prs.slide_width.pt - slide.shapes[0].width.pt) / 2)

        title_font(slide, font_name, 44, 1)

        insert_two_lines(slide, slide.shapes[0].width, slide.shapes[0].height)

        sp = slide.shapes[1].element
        sp.getparent().remove(sp)

    # 본문 슬라이드
    def create_title_and_content_slide(title, content):
        layout = prs.slide_layouts[l_type["title_and_content"]]

        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content

        title_font(slide, font_name, 44)
        text_font(slide, 1, font_name, 20)

    # 이미지가 들어가는 슬라이드
    def create_title_and_content_and_image_slide(title, content):
        layout = prs.slide_layouts[l_type["two_content"]]
    
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
    
        slide.shapes.title.text_frame.word_wrap = True
        title_font(slide, font_name, 44)
        text_font(slide, 1, font_name, 20)
        
        img_path = image_crawler(slide.placeholders[1].text_frame.paragraphs[0].text)

        pic = slide.shapes.add_picture(img_path, slide.placeholders[1].left + slide.placeholders[1].width, slide.placeholders[2].top,
                                    slide.placeholders[2].width)
        pic.top = Pt((prs.slide_height.pt - pic.height.pt) / 2)

        if pic.height > pic.width:
            sp = slide.shapes[3]._element
            sp.getparent().remove(sp)

            pic = slide.shapes.add_picture(img_path, slide.placeholders[1].left + slide.placeholders[1].width, slide.placeholders[2].top,
                                        None, slide.placeholders[2].height)
            pic.left = Pt((prs.slide_width.pt / 2) + ((slide.placeholders[2].width.pt - pic.width.pt) / 2))
        
        sp = slide.shapes[2].element
        sp.getparent().remove(sp)
        
    # ppt에 내용 배치
    def parse_text(lines):
        list_of_slides = "".join(lines).split("[SLIDEBREAK]")

        for slide in list_of_slides:
            slide_type = search_for_slide_type(slide)

            if slide_type == "[L_TS]":
                create_title_slide(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]"),
                                find_text_in_between_tags(str(slide), "[SUBTITLE]", "[/SUBTITLE]"))
                create_index(lines)
            elif slide_type == "[L_CS]":
                create_title_and_content_slide(
                    "".join(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")),
                    "".join(find_text_in_between_tags(str(slide), "[CONTENT]",
                                                        "[/CONTENT]")))
            elif slide_type == "[L_IS]":
                create_title_and_content_and_image_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]",
                                                                                            "[/TITLE]")),
                                                        "".join(find_text_in_between_tags(str(slide), "[CONTENT]",
                                                                                            "[/CONTENT]")))
                                                            #"".join(find_text_in_between_tags(str(slide), "[IMAGE]",
                                                            #                                 "[/IMAGE]"))
            elif slide_type == "[L_THS]":
                create_section_header_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")))

        create_last_slide(last_text)

    # PPT 제목 추출 -> 파일명
    def find_title():
        return prs.slides[0].shapes.title.text

    ### PPT 생성 후 저장 ###
    parse_text(extract_txt)

    # 특수문자가 있으면 파일명 저장시 오류
    save_ppt()
    print("PPT 생성 완료")

    delete_all_slides()

### PPT Generator 모듈의 main 함수 역할 ###
def generator_main(extract_txt):
    global unique_image_name
    unique_image_name = ''.join(random.choice(string.ascii_uppercase + string.ascii_lowercase + string.digits) for _ in
                                range(16))

    origin_text = get_original_text(extract_txt)
    generate_ppt(origin_text)