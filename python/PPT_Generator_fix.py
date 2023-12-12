import base64
import glob
import os
import random
import re
import string
from urllib.parse import urlparse

from icrawler import ImageDownloader
from icrawler.builtin import GoogleImageCrawler

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

import setting

IMAGE_DIR = setting.IMAGE_DIR
RESULT_DIR = setting.RESULT_DIR

# 사용자 지정 옵션
user_path = os.getcwd()
user_data = "NewsText.txt"
last_text = "Q&A"
font_name = "NanumGothic"

ppt_name = []

# 레이아웃 타입
l_type = {"title_slide": 0, 
          "title_and_content": 1, 
          "section_header": 2, 
          "two_content": 3, 
          "comparison": 4, 
          "title_only": 5,
          "blank": 6, 
          "content_with_caption": 7, 
          "picture_with_caption": 8
          }

unique_image_name = None

### 텍스트 및 이미지 관련 ###
# Base64 인코더로 파일명 생성 후 다운로드
class Base64NameDownloader(ImageDownloader):
    def get_filename(self, task, default_ext):
        url_path = urlparse(task['file_url'])[2]
        if '.' in url_path:
            extension = url_path.split('.')[-1]
            if extension.lower() not in ['jpg', 'jpeg', 'png', 'bmp', 'tiff', 'gif', 'ppm', 'pgm']:
                extension = default_ext
        else:
            extension = default_ext

        filename = base64.b64encode(url_path.encode()).decode()
        return "p_" + unique_image_name + '{}.{}'.format(filename, extension)

# 원본 텍스트 불러오기
def get_original_text(text):
    with open(text, 'r', encoding='UTF8') as file:
        lines = file.readlines()
        return lines

# 텍스트 파일 내 태그 제거 및 분류
def find_text_in_between_tags(text, start_tag, end_tag):
    start_pos = text.find(start_tag)
    end_pos = text.find(end_tag)
    result = []

    while start_pos > -1 and end_pos > -1:
        text_between_tags = text[start_pos + len(start_tag):end_pos]
        result.append(text_between_tags)
        start_pos = text.find(start_tag, end_pos + len(end_tag))
        end_pos = text.find(end_tag, start_pos)

    res1 = "".join(result).replace('- ', '')
    res2 = re.sub(r"\[IMAGE\].*?\[/IMAGE\]", '', res1)
    res2 = re.sub(r"\n", '', res2, count=1)

    if len(result) > 0:
        return res2
    else:
        return ""

# 슬라이드 타입 판별
def search_for_slide_type(text):
    tags = ["[L_TS]", "[L_CS]", "[L_IS]", "[L_THS]"]
    found_text = next((s for s in tags if s in text), None)
    return found_text

# 폰트 조정(제목)
def title_font(slide, font_name, font_size, is_bold=None):
    for paragraph in slide.shapes.title.text_frame.paragraphs:
        paragraph.font.name = font_name
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = is_bold

# 폰트 조정(본문)
def text_font(slide, idx, font_name, font_size, is_bold=None):
    for paragraph in slide.placeholders[idx].text_frame.paragraphs:
        paragraph.font.name = font_name
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = is_bold

# 다운로드한 이미지 파일명에 넣을 무작위 문자열
def refresh_unique_image_name():
    global unique_image_name
    unique_image_name = ''.join(random.choice(string.ascii_uppercase + string.ascii_lowercase + string.digits)
                                for _ in range(16))
    return

# 이미지 크롤링
def image_crawler(keyword):
    refresh_unique_image_name()

    google_crawler = GoogleImageCrawler(downloader_cls=Base64NameDownloader, storage={'root_dir': user_path})
    google_crawler.crawl(keyword=keyword, max_num=1)

    dir_path = os.path.dirname(os.path.realpath(__file__))
    file_name = glob.glob(f"p_{unique_image_name}*")

    img_path = os.path.join(dir_path, file_name[0])
    print("file_name: ", file_name[0])
    print("img_path: ", img_path)
    return img_path

# 위아래 선 긋기
def insert_two_lines(slide, shape_width, shape_height):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, slide.shapes[0].left, slide.shapes[0].top, shape_width, 4)
    line.fill.background()
    line.line.color.rgb = RGBColor(0, 0, 0)
    line.line.width = Pt(2)

    line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, slide.shapes[0].left, slide.shapes[0].top + shape_height, shape_width, 4)
    line2.fill.background()
    line2.line.color.rgb = RGBColor(0, 0, 0)
    line2.line.width = Pt(2)

### PPT 생성 모듈 ###
def generate_ppt(filename, lines): 
    prs = Presentation("theme0.pptx")

    ### 슬라이드 제작 함수 ###
    # 모든 슬라이드 삭제
    def delete_all_slides():
        for i in range(len(prs.slides) - 1, -1, -1):
            r_id = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(r_id)
            del prs.slides._sldIdLst[i]

    # 마지막 슬라이드
    def create_last_slide(last_text):
        layout = prs.slide_layouts[l_type["title_only"]]
        last_slide = prs.slides.add_slide(layout)

        width = Inches(6)
        height = Inches(2)

        last_slide.shapes.title.text = last_text
        last_slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        last_slide.shapes[0].width = width
        last_slide.shapes[0].height = height

        last_slide.shapes[0].top = Pt((prs.slide_height.pt - last_slide.shapes[0].height.pt) / 2)
        last_slide.shapes[0].left = Pt((prs.slide_width.pt - last_slide.shapes[0].width.pt) / 2)

        title_font(last_slide, font_name, 44, 1)

        insert_two_lines(last_slide, last_slide.shapes[0].width, last_slide.shapes[0].height)

    # 타이틀 슬라이드
    def create_title_slide(title, subtitle):
        title_slide = prs.slides[0]

        title_slide.shapes.title.text = title
        title_slide.placeholders[1].text = subtitle

        title_font(title_slide, font_name, 56, 1)
        text_font(title_slide, 1, font_name, 30)

        insert_two_lines(title_slide, title_slide.shapes[0].width, title_slide.shapes[0].top + title_slide.shapes[0].height)

        '''
        layout = prs.slide_layouts[l_type["title_slide"]]

        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle
        '''

    # 섹션 헤더 슬라이드
    def create_section_header_slide(title):
        layout = prs.slide_layouts[l_type["section_header"]]

        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        slide.shapes.title.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        width = Inches(6)
        height = Inches(2)

        slide.shapes[0].width = width
        slide.shapes[0].height = height

        slide.shapes[0].top = Pt((prs.slide_height.pt - slide.shapes[0].height.pt) / 2)
        slide.shapes[0].left = Pt((prs.slide_width.pt - slide.shapes[0].width.pt) / 2)

        title_font(slide, font_name, 44, 1)

        insert_two_lines(slide, slide.shapes[0].width, slide.shapes[0].height)

        sp = slide.shapes[1].element
        sp.getparent().remove(sp)

    # 본문 슬라이드
    def create_title_and_content_slide(title, content):
        layout = prs.slide_layouts[l_type["title_and_content"]]

        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content

        title_font(slide, font_name, 44)
        text_font(slide, 1, font_name, 20)

    # 이미지가 들어가는 슬라이드
    def create_title_and_content_and_image_slide(title, content):
        layout = prs.slide_layouts[l_type["two_content"]]
    
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
    
        slide.shapes.title.text_frame.word_wrap = True
        title_font(slide, font_name, 44)
        text_font(slide, 1, font_name, 20)

        #left = top = Cm(3)
        #pic = slide.shapes.add_picture(f"{user_path}/test.jpg", slide.placeholders[1].left+slide.placeholders[1].width,
        #                                   slide.placeholders[2].top, slide.placeholders[2].width)
        #pic.top = Pt((prs.slide_height.pt - pic.height.pt) / 2)
        
        img_path = image_crawler(title)

        pic = slide.shapes.add_picture(img_path, slide.placeholders[1].left + slide.placeholders[1].width, slide.placeholders[2].top,
                                       slide.placeholders[2].width)
        pic.top = Pt((prs.slide_height.pt - pic.height.pt) / 2)

        if pic.height > pic.width:
            sp = slide.shapes[3]._element
            sp.getparent().remove(sp)

            pic = slide.shapes.add_picture(img_path, slide.placeholders[1].left + slide.placeholders[1].width, slide.placeholders[2].top,
                                           None, slide.placeholders[2].height)
            pic.left = Pt((prs.slide_width.pt / 2) + ((slide.placeholders[2].width.pt - pic.width.pt) / 2))
        
        slide.shapes[3].shadow

        sp = slide.shapes[2].element
        sp.getparent().remove(sp)
        
    # ppt에 내용 배치
    def parse_text(lines):
        list_of_slides = "".join(lines).split("[SLIDEBREAK]")

        for slide in list_of_slides:
            slide_type = search_for_slide_type(slide)

            if slide_type == "[L_TS]":
                create_title_slide(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]"),
                                   find_text_in_between_tags(str(slide), "[SUBTITLE]", "[/SUBTITLE]"))
            elif slide_type == "[L_CS]":
                create_title_and_content_slide(
                    "".join(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")),
                    "".join(find_text_in_between_tags(str(slide), "[CONTENT]",
                                                        "[/CONTENT]")))
            elif slide_type == "[L_IS]":
                create_title_and_content_and_image_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]",
                                                                                               "[/TITLE]")),
                                                         "".join(find_text_in_between_tags(str(slide), "[CONTENT]",
                                                                                               "[/CONTENT]")))
                                                             #"".join(find_text_in_between_tags(str(slide), "[IMAGE]",
                                                             #                                 "[/IMAGE]"))
            elif slide_type == "[L_THS]":
                create_section_header_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")))

        create_last_slide(last_text)

    # PPT 제목 추출 -> 파일명
    def find_title():
        return prs.slides[0].shapes.title.text

    ### PPT 생성 후 저장 ###
    parse_text(lines)

    ppt_name = filename.replace("_pptx.txt", "_ppt.pptx")
    # ppt_name = re.sub(r"[^\uAC00-\uD7A30-9a-zA-Z\s]", "", find_title())
    prs.save(f"{RESULT_DIR}/{ppt_name}")
    print("PPT was successfully saved")

    delete_all_slides()

### main 함수 ###
def generator_main(filename, extract_text):
    global unique_image_name
    unique_image_name = ''.join(random.choice(string.ascii_uppercase + string.ascii_lowercase + string.digits) for _ in
                                range(16))   
    
    text = get_original_text(extract_text)
    generate_ppt(filename, text)