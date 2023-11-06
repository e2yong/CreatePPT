import base64
import glob
import os
import random
import re
import string
from urllib.parse import urlparse

from icrawler import ImageDownloader
from icrawler.builtin import GoogleImageCrawler
from pptx import Presentation
from pptx.util import Pt, Cm, Inches
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData, XyChartData
from tkinter import *
from PIL import ImageTk,Image

unique_image_name = None

user_path = os.getcwd()
user_data = "NewsText.txt"

font_name = "NanumGothic"

### 텍스트 및 이미지 파일관리 함수들 ###
# Base64 인코더로 파일명 생성 후 다운로드
class Base64NameDownloader(ImageDownloader):
    def get_filename(self, task, default_ext):
        url_path = urlparse(task['file_url'])[2]
        if '.' in url_path:
            extension = url_path.split('.')[-1]
            if extension.lower() not in ['jpg', 'jpeg', 'png', 'bmp', 'tiff', 'gif', 'ppm', 'pgm']:
                extension = default_ext
        else:
            extension = default_ext

        filename = base64.b64encode(url_path.encode()).decode()
        return "p_" + unique_image_name + '{}.{}'.format(filename, extension)

# 원본 텍스트 불러오기
def get_original_text():
    with open(f"{user_path}/{user_data}", 'r', encoding='UTF8') as file:
        lines = file.readlines()
        return lines

# 텍스트 파일 내 태그 제거 및 분류
def find_text_in_between_tags(text, start_tag, end_tag):
    start_pos = text.find(start_tag)
    end_pos = text.find(end_tag)
    result = []

    while start_pos > -1 and end_pos > -1:
        text_between_tags = text[start_pos + len(start_tag):end_pos]
        result.append(text_between_tags)
        start_pos = text.find(start_tag, end_pos + len(end_tag))
        end_pos = text.find(end_tag, start_pos)

    res1 = "".join(result).replace('- ', '')
    res2 = re.sub(r"\[IMAGE\].*?\[/IMAGE\]", '', res1)
    res2 = re.sub(r"\n", '', res2, count=1)

    if len(result) > 0:
        return res2
    else:
        return ""

# 슬라이드 타입 판별
def search_for_slide_type(text):
    tags = ["[L_TS]", "[L_CS]", "[L_IS]", "[L_THS]"]
    found_text = next((s for s in tags if s in text), None)
    return found_text

# 폰트 조정
def title_font(slide, font_name, font_size):
    for paragraph in slide.shapes.title.text_frame.paragraphs:
        paragraph.font.name = font_name
        paragraph.font.size = Pt(font_size)

def text_font(slide, idx, font_name, font_size):
    for paragraph in slide.placeholders[idx].text_frame.paragraphs:
        paragraph.font.name = font_name
        paragraph.font.size = Pt(font_size)

# 다운로드한 이미지 파일명에 넣을 무작위 문자열
def refresh_unique_image_name():
    global unique_image_name
    unique_image_name = ''.join(random.choice(string.ascii_uppercase + string.ascii_lowercase + string.digits)
                                for _ in range(16))
    return

# 이미지 크롤링
def image_crawler(keyword):
    refresh_unique_image_name()

    google_crawler = GoogleImageCrawler(downloader_cls=Base64NameDownloader, storage={'root_dir': os.getcwd()})
    google_crawler.crawl(keyword=keyword, max_num=1)

    dir_path = os.path.dirname(os.path.realpath(__file__))
    file_name = glob.glob(f"p_{unique_image_name}*")

    img_path = os.path.join(dir_path, file_name[0])
    return img_path

### PPT 생성 모듈 ###
def generate_ppt(lines):
    prs = Presentation("theme0.pptx")

    ### 슬라이드 제작 함수 ###
    # 모든 슬라이드 삭제
    def delete_all_slides():
        for i in range(len(prs.slides) - 1, -1, -1):
            r_id = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(r_id)
            del prs.slides._sldIdLst[i]

    # 타이틀 슬라이드
    def create_title_slide(title, subtitle):
        layout = prs.slide_layouts[0]
        slide = prs.slides[0]

        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

        slide.shapes.title.text_frame.paragraphs[0].font.bold = True

        title_font(slide, font_name, 56)
        text_font(slide, 1, font_name, 20)

    # 섹션 타이틀 슬라이드
    def create_section_header_slide(title):
        layout = prs.slide_layouts[2]

        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title

    # 본문 슬라이드
    def create_title_and_content_slide(title, content):
        layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content

        title_font(slide, font_name, 44)
        text_font(slide, 1, font_name, 20)

    # 이미지가 들어가는 슬라이드
    def create_title_and_content_and_image_slide(title, content):
        layout = prs.slide_layouts[3]

        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content

        slide.shapes.title.text_frame.word_wrap = True
        title_font(slide, font_name, 44)
        text_font(slide, 1, font_name, 20)

        img_path = image_crawler(title)

        pic = slide.shapes.add_picture(img_path, slide.placeholders[1].left+slide.placeholders[1].width, slide.placeholders[2].top,
                                       slide.placeholders[2].width)
        pic.top = Pt((prs.slide_height.pt - pic.height.pt) / 2)

        if pic.height > pic.width:
            sp = slide.shapes[3]._element
            sp.getparent().remove(sp)

            pic = slide.shapes.add_picture(img_path, slide.placeholders[1].left+slide.placeholders[1].width, slide.placeholders[2].top,
                                           None, slide.placeholders[2].height)
            pic.left = Pt((prs.slide_width.pt / 2) + ((slide.placeholders[2].width.pt - pic.width.pt) / 2))
        
        sp = slide.shapes[2].element
        sp.getparent().remove(sp)
    
    # ppt에 내용 배치
    def parse_text(lines):
        list_of_slides = "".join(lines).split("[SLIDEBREAK]")

        for slide in list_of_slides:
            slide_type = search_for_slide_type(slide)

            if slide_type == "[L_TS]":
                create_title_slide(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]"),
                                   find_text_in_between_tags(str(slide), "[SUBTITLE]", "[/SUBTITLE]"))
            elif slide_type == "[L_CS]":
                create_title_and_content_slide(
                    "".join(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")),
                    "".join(find_text_in_between_tags(str(slide), "[CONTENT]",
                                                        "[/CONTENT]")))
            elif slide_type == "[L_IS]":
                create_title_and_content_and_image_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]",
                                                                                               "[/TITLE]")),
                                                         "".join(find_text_in_between_tags(str(slide), "[CONTENT]",
                                                                                               "[/CONTENT]")))
                                                             #"".join(find_text_in_between_tags(str(slide), "[IMAGE]",
                                                             #                                 "[/IMAGE]"))
            elif slide_type == "[L_THS]":
                create_section_header_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")))
    
    def find_title():
        return prs.slides[0].shapes.title.text

    ### PPT 생성 후 저장 ###
    parse_text(lines)

    ppt_name = re.sub(r"[^\uAC00-\uD7A30-9a-zA-Z\s]", "", find_title())
    prs.save(f"{ppt_name}.pptx")
    print("PPT was successfully saved")

    delete_all_slides()

# main 함수
def main():
    global unique_image_name
    unique_image_name = ''.join(random.choice(string.ascii_uppercase + string.ascii_lowercase + string.digits) for _ in
                                range(16))

    text = get_original_text()
    generate_ppt(text)

if __name__ == '__main__':
    main()